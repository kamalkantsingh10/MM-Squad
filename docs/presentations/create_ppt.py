#!/usr/bin/env python3
"""
Generates whitepaper.pptx from whitepaper content.
One slide per major section, diagrams embedded as PNG images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
TITLE_COL   = RGBColor(0xFF, 0xFF, 0xFF)   # white
BODY_COL    = RGBColor(0xCC, 0xD6, 0xE0)   # light grey-blue
ACCENT      = RGBColor(0xF3, 0x9C, 0x12)   # orange
ACCENT2     = RGBColor(0x3B, 0xAA, 0xD9)   # sky blue
INFOSYS     = RGBColor(0x00, 0x70, 0xC0)   # Infosys blue
BANK        = RGBColor(0x70, 0xAD, 0x47)   # green
DIVIDER     = RGBColor(0xF3, 0x9C, 0x12)   # orange divider line

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR    = os.path.join(SCRIPT_DIR, "whitepaper_images")
OUT_FILE   = os.path.join(SCRIPT_DIR, "whitepaper.pptx")

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, colour=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, colour):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, colour=BODY_COL,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txb


def add_title_bar(slide, title, subtitle=None):
    # Orange divider line at top
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), DIVIDER)
    # Title
    add_textbox(slide, title,
                Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75),
                font_size=32, bold=True, colour=TITLE_COL)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.4),
                    font_size=16, colour=ACCENT)
    # Bottom orange line
    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), DIVIDER)


def add_bullet_block(slide, items, left, top, width, height,
                     font_size=15, bullet="▸ "):
    """items = list of (text, is_highlight, is_bold)"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, highlight, bold) in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = ACCENT if highlight else BODY_COL


def add_image(slide, img_path, left, top, width=None, height=None):
    if not os.path.exists(img_path):
        print(f"  WARNING: image not found: {img_path}")
        return
    if width and height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(img_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        slide.shapes.add_picture(img_path, left, top)


def add_two_col_table(slide, headers, rows, left, top, width, height,
                      font_size=12):
    cols = len(headers)
    table = slide.shapes.add_table(len(rows)+1, cols, left, top, width, height).table
    col_w = width // cols
    for i in range(cols):
        table.columns[i].width = col_w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = ACCENT

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            alt = RGBColor(0x10, 0x24, 0x38) if i % 2 == 0 else RGBColor(0x14, 0x2C, 0x44)
            cell.fill.fore_color.rgb = alt
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font_size)
            run.font.color.rgb = BODY_COL


def img(n):
    return os.path.join(IMG_DIR, f"diagram_{n}.png")


# ── Build slides ──────────────────────────────────────────────────────────────

def build(prs):

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.1), DIVIDER)
    add_rect(s, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), DIVIDER)
    add_rect(s, 0, Inches(2.8), SLIDE_W, Inches(2.0), RGBColor(0x1A, 0x3A, 0x5C))
    add_textbox(s, "Mainframe Modernisation\nwith AI Agents",
                Inches(0.8), Inches(2.9), Inches(11.0), Inches(1.6),
                font_size=44, bold=True, colour=TITLE_COL, align=PP_ALIGN.CENTER)
    add_textbox(s, "A Practical Approach — Powered by Tools You Already Have",
                Inches(0.8), Inches(4.5), Inches(11.0), Inches(0.5),
                font_size=20, colour=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, "Infosys  ·  Confidential  ·  2026",
                Inches(0.8), Inches(6.8), Inches(11.0), Inches(0.4),
                font_size=13, colour=RGBColor(0x66, 0x88, 0xAA), align=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Executive Summary")
    bullets = [
        ("COBOL estates consume 70–75% of IT budgets. 70% of modernisation projects fail with standard agile.", True, True),
        ("Existing AI tools require 3–12 months procurement before work begins — and leave implementation to the client.", False, False),
        ("This proposal uses GitHub Copilot, GitLab and BMAD — already licensed. Zero new software. Work begins next sprint.", True, True),
        ("3-stage engagement: 6-month pilot → 7-month full application → ongoing factory at 2.5 HC per app/year.", False, False),
    ]
    add_bullet_block(s, bullets, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), font_size=17)

    # ── Slide 3: The Problem ──────────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "The Problem", "Why existing tools don't work for a regulated bank")
    bullets = [
        ("Months before work begins — vendor approval in regulated IT takes 3–12 months before technical setup even starts.", True, True),
        ("$2M–$15M+ absorbed by licensing, server infrastructure, and LLM hosting — before delivery begins.", True, True),
        ("Locked to yesterday's models — IBM WCA for Z runs Mistral / Granite (fixed). Cannot access frontier models without rebuilding infrastructure.", False, True),
        ("Products deliver tooling, not outcomes — end-to-end implementation left entirely to the client.", False, False),
        ("70% of projects fail — standard agile assumes you understand what you're building. COBOL breaks that assumption.", False, False),
    ]
    add_bullet_block(s, bullets, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), font_size=16)

    # ── Slide 4: What the Organisation Already Has ────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "What You Already Have", "The tools for this solution are already licensed and running")

    boxes = [
        ("GitHub Copilot", "42% market share · 20M users · 90% of Fortune 100\nClaude Opus 4.6 + GPT-5.3-Codex available today\nModels refreshed continuously — deprecated within days, not years\nNo additional cost to any licensed developer", ACCENT2),
        ("GitHub Copilot Agent Mode", "Full agentic framework — plans, builds, reviews, delegates\nCustom agents via markdown files\nMCP servers connect without additional infrastructure\nDevelopers already using it — zero learning curve", ACCENT),
        ("GitLab", "End-to-end software lifecycle platform\nNative AI: CI/CD, issue generation, test creation, code review\nMaps directly to audit and traceability requirements\nEvery pipeline stage tracks in GitLab — zero additional tooling", BANK),
    ]
    box_w = Inches(4.1)
    for i, (title, body, col) in enumerate(boxes):
        left = Inches(0.3 + i * 4.35)
        add_rect(s, left, Inches(1.3), box_w, Inches(0.45), col)
        add_textbox(s, title, left + Inches(0.1), Inches(1.32), box_w - Inches(0.2), Inches(0.4),
                    font_size=15, bold=True, colour=RGBColor(0x0D, 0x1B, 0x2A))
        add_textbox(s, body, left + Inches(0.1), Inches(1.85), box_w - Inches(0.2), Inches(5.0),
                    font_size=13, colour=BODY_COL)

    # ── Slide 5: Solution ─────────────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Solution — A Mainframe Modernisation Factory")
    bullets = [
        ("Zero procurement — built on GitHub Copilot, GitLab, and BMAD. Technical prerequisites already in place. Work begins this sprint.", True, True),
        ("Complete SDLC — structural analysis → dependency mapping → business rule extraction → architecture → code generation → QA.", False, False),
        ("Frontier AI — Claude Opus 4.6 + GPT-5.3-Codex via GitHub Copilot, refreshed continuously. Customised to org-specific macros and glossary.", True, True),
        ("Documentation by design — every workflow produces business rules, dependency maps, and a queryable spec layer as direct output.", False, False),
        ("Maintenance-ready — analysis workflows apply to impact assessment, developer onboarding, and pre-release verification.", False, False),
        ("Tool agnostic — compatible with any MCP-supported coding assistant. Not locked to any single model or provider.", False, False),
    ]
    add_bullet_block(s, bullets, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), font_size=16)

    # ── Slide 6: How It Works — Three Pillars ────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "How It Works — Three Pillars")
    add_image(s, img(1), Inches(0.3), Inches(1.3), height=Inches(3.8))
    bullets = [
        ("Agents execute — specialised AI for each stage: analysis, architecture, delivery, code generation, QA.", False, True),
        ("Process governs — every stage is a defined workflow with explicit inputs, outputs, and a GitLab review gate.", False, True),
        ("People decide — analyst, architect, domain expert, tech lead and developers approve every output before it moves downstream.", True, True),
        ("Spec layer (SQLite) is the shared backbone — the only shared state between agents. No direct agent-to-agent communication.", False, False),
    ]
    add_bullet_block(s, bullets, Inches(5.5), Inches(1.4), Inches(7.5), Inches(5.5), font_size=15)

    # ── Slide 7: The Agent Squad ──────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "The Agent Squad", "Each agent has one role. Every action is logged to GitLab.")
    add_image(s, img(2), Inches(0.3), Inches(1.3), height=Inches(2.8))
    agents = [
        ("Po", "COBOL analysis: structure, dependencies, business rules → spec layer. Most critical agent. Standalone-capable.", True),
        ("Oogway", "Migration architect — reads spec layer, produces target architecture and migration sequence.", False),
        ("Shifu", "Delivery manager — GitLab project, sprints, milestones, live dashboard.", False),
        ("Tigress / Viper / Monkey", "Code generation — Java, COBOL refactor, Python. Target language set by Oogway.", False),
        ("Mantis", "QA validation — generated code vs spec layer business rules. Manages GitLab epic sign-off.", False),
        ("Group Mode", "All agents active simultaneously — for production incidents, critical decisions, or unblocking sessions.", True),
    ]
    txb = s.shapes.add_textbox(Inches(5.5), Inches(1.35), Inches(7.5), Inches(5.8))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for name, desc, hl in agents:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = f"▸ {name}  "
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT if hl else ACCENT2
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = BODY_COL

    # ── Slide 8: The Process ──────────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "The Process", "Every review gate is handled directly in GitLab")
    add_image(s, img(4), Inches(0.2), Inches(1.2), height=Inches(3.5))
    add_image(s, img(3), Inches(6.8), Inches(1.2), height=Inches(3.5))
    stages = [
        ("Inception (2w)",     "App Owner + Architect",                    "Project brief, tech stack, success criteria"),
        ("Setup (1w)",         "Technical Lead",                           "Agent config, GitLab init, spec layer setup"),
        ("Plan (2w)",          "App Owner + Architect + Domain Expert",    "PRD, epics, architecture, testing strategy"),
        ("Dev (~1w/epic)",     "BA + App Team Tech Lead",                  "Business rules → user stories → code. Parallel epics."),
        ("Testing",            "App Owner + App Team Tech Lead",           "Mantis validation + UAT. Behaviour must match or deviation documented."),
        ("Production (1m)",    "App Owner",                                "Monitored against Inception success criteria"),
        ("Maintenance",        "Incident Owner",                           "Group Mode or targeted agents. Same GitLab-tracked pattern."),
    ]
    add_two_col_table(s,
        ["Stage", "Sign-off", "Output"],
        stages,
        Inches(0.2), Inches(4.85), Inches(12.9), Inches(2.5),
        font_size=10)

    # ── Slide 9: People ───────────────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "People", "Agents execute. People decide. No new headcount required.")

    rows = [
        ("App Owner",            "🏦 Bank",    "✅ Approve",              "Final accountability · all sign-offs",                    "All"),
        ("App Team Tech Lead",   "🏦 Bank",    "✅ Approve",              "Technical quality gate · dev + test sign-off",            "Setup, Dev, Testing"),
        ("Business Domain Expert","🏦 Bank",   "🤝 Support",              "Business context · validates extracted business rules",    "Plan, Dev"),
        ("Delivery Lead",        "🔵 Infosys", "⚙️ Execute",              "Runs engagement · operates Shifu · GitLab sprints",       "All"),
        ("COBOL Analyst",        "🔵 Infosys", "⚙️ Execute",              "Drives Po workflows · reviews analysis output",           "Setup, Dev"),
        ("Business Analyst",     "🔵 Infosys", "⚙️ Execute",              "User stories from business rules · epic backlog",         "Dev"),
        ("Developers",           "🔵 Infosys", "⚙️ Execute",              "Review, refine, and own all generated code",              "Dev"),
        ("Enterprise Architect", "🔵 Infosys", "⚙️ Execute + ✅ Approve", "Produces + approves migration architecture",              "Inception, Plan"),
        ("Project Manager",      "🔵 Infosys", "⚙️ Execute",              "Governance · risk · budget · stakeholder comms",          "All"),
    ]
    add_two_col_table(s,
        ["Role", "By", "Mode", "Responsibility", "Stages"],
        rows,
        Inches(0.2), Inches(1.3), Inches(12.9), Inches(5.6),
        font_size=11)

    add_textbox(s,
        "For ~700 modules / 500,000 LOC: estimated 3–5 annualised Infosys HC for full engagement",
        Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
        font_size=12, colour=ACCENT, bold=True)

    # ── Slide 10: Building the Factory ───────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Building the Factory", "Incrementally — each stage delivers output and builds capability")
    add_image(s, img(5), Inches(0.2), Inches(1.3), width=Inches(8.0))

    stages_f = [
        ("Stage 1 — Pilot", "6 months · 6 HC (3 EU + 3 India)", "Build agent v1 · align process · prove pipeline on one epic"),
        ("Stage 2 — First App", "7 months · 6 HC (same team)", "Full application modernised · agent v2 refined from real-world data"),
        ("Stage 3 — Factory", "Ongoing · 2.5 HC per app/year", "Repeatable at scale · team sized by annual application throughput"),
    ]
    for i, (title, meta, desc) in enumerate(stages_f):
        top = Inches(1.5 + i * 1.8)
        add_rect(s, Inches(8.5), top, Inches(0.06), Inches(1.4), ACCENT)
        add_textbox(s, title, Inches(8.7), top, Inches(4.4), Inches(0.45),
                    font_size=14, bold=True, colour=ACCENT)
        add_textbox(s, meta, Inches(8.7), top + Inches(0.42), Inches(4.4), Inches(0.35),
                    font_size=12, colour=ACCENT2)
        add_textbox(s, desc, Inches(8.7), top + Inches(0.78), Inches(4.4), Inches(0.55),
                    font_size=12, colour=BODY_COL)

    # ── Slide 11: Appendix — Market Context ──────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Appendix A — Market Context")
    add_image(s, img(6), Inches(0.2), Inches(1.3), height=Inches(3.5))
    stats = [
        ("Market size 2025",                "$8.4B (mainframe modernisation)"),
        ("Projected 2030",                  "$13.3B  ·  CAGR 9.7%"),
        ("Daily commerce on COBOL",         "$3 trillion"),
        ("Avg COBOL programmer age",        "~60 years · one-third retiring by 2030"),
        ("IT budget on legacy maintenance", "70–75% (financial services)"),
        ("Project failure rate",            "70% with standard agile"),
        ("IBM stock drop (Feb 2026)",        "–13% in one day on Anthropic COBOL announcement"),
        ("IBM market cap lost",             "$31 billion in one day"),
    ]
    add_two_col_table(s,
        ["Statistic", "Figure"],
        stats,
        Inches(6.5), Inches(1.3), Inches(6.6), Inches(5.8),
        font_size=11)

    # ── Slide 12: Appendix — Competitive Comparison ───────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Appendix B — Competitive Comparison")
    rows_c = [
        ("IBM WCA for Z",          "Months + licensing",          "Mistral / Granite (fixed)", "✗", "✗", "✗"),
        ("AWS Mainframe Mod.",      "Months + cloud setup",        "AWS-hosted (fixed)",        "✗", "✗", "✗"),
        ("Micro Focus / OpenText",  "Months + licensing",          "Proprietary (fixed)",       "✗", "✗", "✗"),
        ("Microsoft / Bankdata",    "Setup required",              "Open source (fixed)",       "✗", "✗", "✗"),
        ("This proposal",           "Zero — start next sprint",    "Claude Opus 4.6 + GPT-5.3-Codex (live)", "✅", "✅", "✅"),
    ]
    add_two_col_table(s,
        ["Tool", "Procurement", "AI Model", "Business Rules", "Structured IR", "End-to-End"],
        rows_c,
        Inches(0.2), Inches(1.4), Inches(12.9), Inches(5.6),
        font_size=12)

    # ── Slide 13: Appendix — Security ────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Appendix C — Security & Data Privacy", "Designed local-first for regulated environments")
    add_image(s, img(7), Inches(0.2), Inches(1.3), height=Inches(3.5))
    sec = [
        ("Source code stays local",        "All MCP servers run locally via STDIO — no source transmitted externally"),
        ("Spec layer stays local",          "SQLite on local machine · configurable path · gitignored"),
        ("LLM context is scoped",           "Only the specific paragraph cluster under analysis — no full program dumps"),
        ("GitLab receives metadata only",   "Stage labels, status comments, summaries — not raw COBOL source"),
        ("Credentials via env vars",        "GITLAB_TOKEN in environment — never written to config files"),
    ]
    add_two_col_table(s,
        ["Principle", "Implementation"],
        sec,
        Inches(6.0), Inches(1.4), Inches(7.1), Inches(5.6),
        font_size=11)

    # ── Slide 14: Appendix — Agent Capabilities ───────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    add_title_bar(s, "Appendix D — Agent Capabilities Matrix")
    add_image(s, img(8), Inches(0.2), Inches(1.3), height=Inches(5.8))
    matrix = [
        ("Po",                       "✅", "✅", "✅", "✅", "✅"),
        ("Oogway",                   "—",  "✅", "—",  "—",  "✅"),
        ("Shifu",                    "—",  "✅", "—",  "—",  "✅"),
        ("Tigress / Viper / Monkey", "✅", "✅", "✅", "—",  "✅"),
        ("Mantis",                   "—",  "✅", "—",  "—",  "✅"),
    ]
    add_two_col_table(s,
        ["Agent", "cobol-parser", "specdb", "delta-macros", "jcl-parser", "gitlab"],
        matrix,
        Inches(7.0), Inches(1.4), Inches(6.1), Inches(3.8),
        font_size=12)


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = new_prs()
    build(prs)
    prs.save(OUT_FILE)
    print(f"✓ Saved: {OUT_FILE}  ({len(prs.slides)} slides)")
